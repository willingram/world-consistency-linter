from __future__ import annotations

import csv
import zipfile
from typing import Any

from .models import ExtractionResult, FileSpec, SourceRef, TextChunk

TEXT_EXTS = {".txt", ".md", ".json", ".yaml", ".yml"}
CSV_EXTS = {".csv", ".tsv"}


def extract_all(files: list[FileSpec]) -> ExtractionResult:
    chunks: list[TextChunk] = []
    errors: list[str] = []
    image_only: list[str] = []
    extractor_disagreements: list[str] = []
    for spec in files:
        file_chunks, file_errors, file_image_only, file_disagreements = extract_file(spec)
        chunks.extend(file_chunks)
        errors.extend(file_errors)
        extractor_disagreements.extend(file_disagreements)
        if file_image_only:
            image_only.append(spec.path)
        if not _has_content_chunks(file_chunks) and not file_image_only and not file_errors:
            errors.append(f"{spec.path}: empty extraction")
    return ExtractionResult(chunks, errors, image_only, extractor_disagreements)


def extract_file(spec: FileSpec) -> tuple[list[TextChunk], list[str], bool, list[str]]:
    chunks = [TextChunk(SourceRef(spec.path, "filename", "path"), spec.path)]
    errors: list[str] = []
    disagreements: list[str] = []
    image_only = False
    if not spec.full_path.exists():
        return chunks, [f"{spec.path}: missing file"], False, disagreements

    try:
        suffix = spec.full_path.suffix.lower()
        if suffix == ".pdf":
            body_chunks, image_only, disagreements = _extract_pdf(spec)
        elif suffix == ".docx":
            body_chunks = _extract_docx(spec)
        elif suffix in {".xlsx", ".xlsm", ".xls"}:
            body_chunks = _extract_xlsx(spec)
        elif suffix == ".pptx":
            body_chunks = _extract_pptx(spec)
        elif suffix in CSV_EXTS:
            body_chunks = _extract_csv(spec)
        elif suffix in TEXT_EXTS:
            body_chunks = _extract_text(spec)
        else:
            body_chunks = _extract_text(spec)
        chunks.extend(body_chunks)
    except Exception as exc:  # noqa: BLE001 - parser exception classes vary by backend
        errors.append(f"{spec.path}: unreadable ({type(exc).__name__}: {exc})")

    chunks.extend(_metadata_chunks(spec))
    return chunks, errors, image_only, disagreements


def _extract_text(spec: FileSpec) -> list[TextChunk]:
    text = spec.full_path.read_text(encoding="utf-8", errors="replace")
    return [TextChunk(SourceRef(spec.path, "text", "whole file"), text)] if text.strip() else []


def _extract_csv(spec: FileSpec) -> list[TextChunk]:
    delimiter = "\t" if spec.full_path.suffix.lower() == ".tsv" else ","
    rows: list[str] = []
    with spec.full_path.open("r", encoding="utf-8-sig", errors="replace", newline="") as handle:
        reader = csv.reader(handle, delimiter=delimiter)
        for index, row in enumerate(reader, start=1):
            rows.append(f"row {index}: " + " | ".join(row))
    text = "\n".join(rows)
    return [TextChunk(SourceRef(spec.path, "csv", "rows"), text)] if text.strip() else []


def _extract_pdf(spec: FileSpec) -> tuple[list[TextChunk], bool, list[str]]:
    import pdfplumber
    from pypdf import PdfReader

    chunks: list[TextChunk] = []
    plumber_text = []
    pypdf_text = []

    with pdfplumber.open(spec.full_path) as pdf:
        for page_num, page in enumerate(pdf.pages, start=1):
            text = page.extract_text() or ""
            if text.strip():
                chunks.append(TextChunk(SourceRef(spec.path, "pdfplumber", f"page {page_num}"), text))
                plumber_text.append(text)

    reader = PdfReader(str(spec.full_path))
    for page_num, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        if text.strip():
            chunks.append(TextChunk(SourceRef(spec.path, "pypdf", f"page {page_num}"), text))
            pypdf_text.append(text)

    if reader.metadata:
        meta_text = "\n".join(f"{key}: {value}" for key, value in reader.metadata.items() if value)
        if meta_text.strip():
            chunks.append(TextChunk(SourceRef(spec.path, "metadata", "pdf info"), meta_text))

    disagreements = []
    if plumber_text and pypdf_text and _normalise_text("\n".join(plumber_text)) != _normalise_text("\n".join(pypdf_text)):
        disagreements.append(f"{spec.path}: pdfplumber and pypdf extracted materially different text")
    return chunks, not plumber_text and not pypdf_text, disagreements


def _extract_docx(spec: FileSpec) -> list[TextChunk]:
    import docx

    doc = docx.Document(str(spec.full_path))
    chunks: list[TextChunk] = []
    body = "\n".join(paragraph.text for paragraph in doc.paragraphs if paragraph.text.strip())
    if body.strip():
        chunks.append(TextChunk(SourceRef(spec.path, "docx-body", "paragraphs"), body))

    table_rows: list[str] = []
    for table_idx, table in enumerate(doc.tables, start=1):
        for row_idx, row in enumerate(table.rows, start=1):
            table_rows.append(f"table {table_idx} row {row_idx}: " + " | ".join(cell.text.strip() for cell in row.cells))
    if table_rows:
        chunks.append(TextChunk(SourceRef(spec.path, "docx-body", "tables"), "\n".join(table_rows)))

    header_footer_text: list[str] = []
    for section_idx, section in enumerate(doc.sections, start=1):
        for label, container in (("header", section.header), ("footer", section.footer)):
            text = "\n".join(p.text for p in container.paragraphs if p.text.strip())
            if text.strip():
                header_footer_text.append(f"section {section_idx} {label}: {text}")
    if header_footer_text:
        chunks.append(TextChunk(SourceRef(spec.path, "docx-header-footer", "sections"), "\n".join(header_footer_text)))

    chunks.extend(_extract_docx_zip_parts(spec, {"comments", "footnotes", "endnotes"}))
    return chunks


def _extract_docx_zip_parts(spec: FileSpec, wanted: set[str]) -> list[TextChunk]:
    chunks: list[TextChunk] = []
    with zipfile.ZipFile(spec.full_path) as archive:
        for name in archive.namelist():
            lower = name.lower()
            if not any(part in lower for part in wanted):
                continue
            raw = archive.read(name).decode("utf-8", errors="ignore")
            text = _strip_xml(raw)
            if text.strip():
                chunks.append(TextChunk(SourceRef(spec.path, "docx-xml", name), text))
    return chunks


def _extract_xlsx(spec: FileSpec) -> list[TextChunk]:
    from openpyxl import load_workbook

    chunks: list[TextChunk] = []
    workbook = load_workbook(spec.full_path, data_only=False, read_only=False)
    chunks.append(TextChunk(SourceRef(spec.path, "xlsx", "sheet names"), ", ".join(workbook.sheetnames)))

    for sheet in workbook.worksheets:
        rows: list[str] = []
        state = "hidden" if sheet.sheet_state != "visible" else "visible"
        for row in sheet.iter_rows():
            values = [_cell_text(cell.value) for cell in row]
            if any(values):
                rows.append(" | ".join(values))
        if rows:
            chunks.append(TextChunk(SourceRef(spec.path, f"xlsx-{state}", sheet.title), "\n".join(rows)))

    defined_names = []
    for name in workbook.defined_names.values():
        defined_names.append(f"{name.name}: {name.attr_text}")
    if defined_names:
        chunks.append(TextChunk(SourceRef(spec.path, "xlsx", "defined names"), "\n".join(defined_names)))

    props = workbook.properties
    prop_text = "\n".join(
        f"{key}: {value}"
        for key, value in {
            "creator": props.creator,
            "lastModifiedBy": props.lastModifiedBy,
            "created": props.created,
            "modified": props.modified,
        }.items()
        if value
    )
    if prop_text:
        chunks.append(TextChunk(SourceRef(spec.path, "metadata", "workbook properties"), prop_text))
    return chunks


def _extract_pptx(spec: FileSpec) -> list[TextChunk]:
    from pptx import Presentation

    prs = Presentation(str(spec.full_path))
    chunks: list[TextChunk] = []
    for idx, slide in enumerate(prs.slides, start=1):
        texts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text)
        if texts:
            chunks.append(TextChunk(SourceRef(spec.path, "pptx-slide", f"slide {idx}"), "\n".join(texts)))
        if slide.has_notes_slide:
            notes = []
            for shape in slide.notes_slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    notes.append(shape.text)
            if notes:
                chunks.append(TextChunk(SourceRef(spec.path, "pptx-notes", f"slide {idx}"), "\n".join(notes)))
    return chunks


def _metadata_chunks(spec: FileSpec) -> list[TextChunk]:
    values: dict[str, Any] = {
        "purported_author": spec.purported_author,
        "purported_org": spec.purported_org,
        "purported_date": spec.purported_date,
    }
    text = "\n".join(f"{key}: {value}" for key, value in values.items() if value is not None)
    return [TextChunk(SourceRef(spec.path, "manifest", "file declaration"), text)] if text else []


def _cell_text(value: Any) -> str:
    if value is None:
        return ""
    return str(value)


def _strip_xml(raw: str) -> str:
    import re

    return re.sub(r"<[^>]+>", " ", raw).replace("&amp;", "&").replace("&lt;", "<").replace("&gt;", ">")


def _has_content_chunks(chunks: list[TextChunk]) -> bool:
    return any(chunk.text.strip() and chunk.source.channel not in {"filename", "manifest"} for chunk in chunks)


def _normalise_text(value: str) -> str:
    import re

    return re.sub(r"\s+", " ", value).strip()
